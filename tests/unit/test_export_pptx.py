"""Unit tests — PPTX export (issue #8)."""

from __future__ import annotations

import stat

import pytest
from pptx import Presentation

from quip_export.formats.pptx import export_pptx

HTML_SINGLE_SLIDE = """
<html><body>
<div class="slide"><h1>Slide Title</h1><p>Body content here.</p></div>
</body></html>
"""

HTML_MULTI_SLIDE = """
<html><body>
<div class="slide"><h1>Intro</h1><p>Welcome</p></div>
<div class="slide"><h1>Agenda</h1><ul><li>Topic A</li><li>Topic B</li></ul></div>
<div class="slide"><h1>Results</h1><p>We did great.</p></div>
</body></html>
"""

HTML_TITLE_ONLY_SLIDE = """
<html><body>
<div class="slide"><h1>Only a title</h1></div>
</body></html>
"""

HTML_NO_HEADING_SLIDE = """
<html><body>
<div class="slide"><p>No heading, just body.</p></div>
</body></html>
"""


class TestExportPptx:
    def test_creates_file(self, tmp_path):
        out = tmp_path / "deck.pptx"
        export_pptx(HTML_SINGLE_SLIDE, out)
        assert out.exists()
        assert out.stat().st_size > 0

    def test_single_slide_creates_one_slide(self, tmp_path):
        out = tmp_path / "deck.pptx"
        export_pptx(HTML_SINGLE_SLIDE, out)
        prs = Presentation(out)
        assert len(prs.slides) == 1

    def test_multi_slide_correct_count(self, tmp_path):
        out = tmp_path / "deck.pptx"
        export_pptx(HTML_MULTI_SLIDE, out)
        prs = Presentation(out)
        assert len(prs.slides) == 3

    def test_slide_title_extracted_from_h1(self, tmp_path):
        out = tmp_path / "deck.pptx"
        export_pptx(HTML_SINGLE_SLIDE, out)
        prs = Presentation(out)
        slide = prs.slides[0]
        title_shape = slide.shapes.title
        assert title_shape is not None
        assert "Slide Title" in title_shape.text

    def test_slide_body_text_present(self, tmp_path):
        out = tmp_path / "deck.pptx"
        export_pptx(HTML_SINGLE_SLIDE, out)
        prs = Presentation(out)
        slide = prs.slides[0]
        all_text = " ".join(shape.text for shape in slide.shapes if shape.has_text_frame)
        assert "Body content" in all_text

    def test_bullet_list_preserved_as_lines(self, tmp_path):
        out = tmp_path / "deck.pptx"
        export_pptx(HTML_MULTI_SLIDE, out)
        prs = Presentation(out)
        agenda_slide = prs.slides[1]
        all_text = " ".join(shape.text for shape in agenda_slide.shapes if shape.has_text_frame)
        assert "Topic A" in all_text
        assert "Topic B" in all_text

    def test_title_only_slide_no_error(self, tmp_path):
        out = tmp_path / "deck.pptx"
        export_pptx(HTML_TITLE_ONLY_SLIDE, out)
        prs = Presentation(out)
        assert len(prs.slides) == 1

    def test_empty_presentation_zero_slides(self, tmp_path, html_empty):
        out = tmp_path / "empty.pptx"
        export_pptx(html_empty, out)
        prs = Presentation(out)
        assert len(prs.slides) == 0

    def test_image_replaced_with_placeholder_text(self, tmp_path, html_with_image):
        out = tmp_path / "img.pptx"
        export_pptx(html_with_image, out)
        prs = Presentation(out)
        all_text = " ".join(
            shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
        )
        assert "[image]" in all_text.lower() or "image" in all_text.lower()

    def test_write_failure_raises_os_error(self, tmp_path):
        ro = tmp_path / "ro"
        ro.mkdir()
        ro.chmod(stat.S_IREAD | stat.S_IEXEC)
        out = ro / "deck.pptx"
        try:
            with pytest.raises((PermissionError, OSError)):
                export_pptx(HTML_SINGLE_SLIDE, out)
        finally:
            ro.chmod(stat.S_IRWXU)
