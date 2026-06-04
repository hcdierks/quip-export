"""Export Quip slides HTML to PPTX."""

from __future__ import annotations

from pathlib import Path

from bs4 import BeautifulSoup, Tag
from pptx import Presentation


def _slide_has_content(body: Tag) -> bool:
    return bool(body.find(["h1", "h2", "h3", "p", "li", "img"]))


def _build_body_lines(container: Tag) -> list[str]:
    lines: list[str] = []
    for elem in container.find_all(["p", "li", "img"]):
        if elem.name == "img":
            lines.append("[image]")
        else:
            text = elem.get_text(strip=True)
            if text:
                lines.append(text)
    return lines


def _add_slide(prs: Presentation, title_text: str | None, body_lines: list[str]) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    title_shape = slide.shapes.title
    if title_shape and title_text:
        title_shape.text = title_text

    if body_lines and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                tf.add_paragraph().text = line


def export_pptx(html: str, output_path: Path) -> None:
    soup = BeautifulSoup(html, "lxml")
    prs = Presentation()

    slide_divs = soup.find_all("div", class_="slide")

    if slide_divs:
        for div in slide_divs:
            h1 = div.find("h1")
            title_text = h1.get_text(strip=True) if h1 else None
            body_lines = _build_body_lines(div)
            _add_slide(prs, title_text, body_lines)
    else:
        body = soup.find("body")
        if body and _slide_has_content(body):  # type: ignore[arg-type]
            h1 = body.find("h1")
            title_text = h1.get_text(strip=True) if h1 else None
            body_lines = _build_body_lines(body)  # type: ignore[arg-type]
            _add_slide(prs, title_text, body_lines)

    prs.save(output_path)
